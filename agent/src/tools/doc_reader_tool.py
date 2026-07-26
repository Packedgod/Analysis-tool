"""Universal document reader: dispatches by file extension.

Supported formats:
  - PDF (.pdf) — pypdfium2 + OCR fallback for image pages
  - Word (.docx) — python-docx (paragraphs + table cells)
  - Excel (.xlsx/.xls) — pandas preview, all sheets
  - PowerPoint (.pptx) — python-pptx (slide text)
  - Images (.png/.jpg/.jpeg/.gif/.bmp/.webp/.tiff) — OCR
  - Plain text (.txt/.md/.log/.json/.yaml/.yml/.toml/.ini/.cfg/.csv/.tsv/
                .html/.xml/.rst/.sql/.sh and common source-code extensions)

All handlers return the same JSON envelope: status, file, format, char_count,
truncated, text. PDF/Excel add format-specific metadata (pages, sheets, ...).
"""

from __future__ import annotations

import contextvars
import json
from pathlib import Path
from typing import Any, Callable

from src.agent.progress import emit_progress
from src.agent.tools import BaseTool
from src.security.scanner import with_security_warnings
from src.tools.path_utils import safe_document_path

# Default read window, and the hard ceiling a single call may ever return so an
# oversized ``max_chars`` cannot overflow the model context. Larger documents are
# paged: each response carries ``next_offset`` to fetch the following window,
# giving full-document access instead of a silent one-shot truncation.
_DEFAULT_MAX_CHARS = 50_000
_HARD_MAX_CHARS = 200_000
_MIN_TEXT_PER_PAGE = 50

# Per-call read window (offset, max_chars), set by read_document and read by the
# shared envelope. A ContextVar keeps it correct under concurrent tool calls.
_window_params: contextvars.ContextVar[tuple[int, int]] = contextvars.ContextVar(
    "doc_reader_window", default=(0, _DEFAULT_MAX_CHARS)
)
_ENCODING_FALLBACK = ("utf-8", "utf-8-sig", "gbk", "gb2312", "big5", "latin-1")

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tiff", ".tif"}
_TEXT_EXTS = {
    # docs / structured
    ".txt", ".md", ".log", ".rst",
    ".json", ".yaml", ".yml", ".toml", ".ini", ".cfg", ".conf", ".env",
    ".csv", ".tsv", ".html", ".htm", ".xml",
    # source code (best-effort, LLM can read raw)
    ".py", ".pyi", ".js", ".jsx", ".ts", ".tsx",
    ".go", ".rs", ".java", ".kt", ".swift",
    ".c", ".h", ".cpp", ".hpp", ".cc",
    ".rb", ".php", ".pl", ".lua",
    ".sh", ".bash", ".zsh", ".ps1", ".bat",
    ".sql", ".r", ".m",
    ".dockerfile", ".makefile", ".cmake",
}

from src.tools.ocr import get_ocr_engine, get_ocr_install_hint

_cached_ocr_engine = None
_cached_ocr_checked = False


def _get_ocr():
    """Return the configured OCR engine (cached), or None."""
    global _cached_ocr_engine, _cached_ocr_checked
    if not _cached_ocr_checked:
        _cached_ocr_engine = get_ocr_engine()
        _cached_ocr_checked = True
    return _cached_ocr_engine


def _ocr_available() -> bool:
    return _get_ocr() is not None


def _ocr_image_array(img) -> str:
    """Run OCR on a numpy image via the pluggable engine."""
    engine = _get_ocr()
    if engine is None:
        return ""
    return engine.recognize(img)


# ---------------- shared helpers ----------------

def _err(msg: str) -> str:
    return json.dumps({"status": "error", "error": msg}, ensure_ascii=False)


def _window(text: str) -> tuple[str, int, int, int | None, bool]:
    """Slice ``text`` to the active read window.

    Returns (slice, offset, returned, next_offset, truncated). ``max_chars`` is
    clamped to ``_HARD_MAX_CHARS``; ``next_offset`` is the offset to request next
    (or None when the window reaches the end).
    """
    offset, max_chars = _window_params.get()
    offset = max(0, int(offset))
    max_chars = min(max(1, int(max_chars)), _HARD_MAX_CHARS)
    total = len(text)
    body = text[offset : offset + max_chars]
    returned = len(body)
    next_offset = offset + returned if offset + returned < total else None
    truncated = next_offset is not None or offset > 0
    return body, offset, returned, next_offset, truncated


def _envelope(path: Path, fmt: str, text: str, **extra: Any) -> str:
    """Build the standard JSON response for the active read window."""
    body, offset, returned, next_offset, truncated = _window(text)
    payload: dict[str, Any] = {
        "status": "ok",
        "file": path.name,
        "format": fmt,
        "char_count": len(text),
        "offset": offset,
        "returned": returned,
        "next_offset": next_offset,
        "truncated": truncated,
        "text": body,
    }
    payload.update(extra)
    payload = with_security_warnings(payload, fields=("text",))
    return json.dumps(payload, ensure_ascii=False)


# ---------------- PDF ----------------

def _parse_pages(pages_str: str, total: int) -> list[int]:
    """Parse '1-10' / '5' / '1,3,5-8' into zero-based indices."""
    out: list[int] = []
    for part in pages_str.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            start, end = part.split("-", 1)
            s = max(int(start.strip()) - 1, 0)
            e = min(int(end.strip()), total)
            out.extend(range(s, e))
        elif part.isdigit():
            out.append(int(part) - 1)
    return sorted(set(out))


def _read_pdf(path: Path, pages: str) -> str:
    """Extract PDF text; OCR pages with too little text."""
    try:
        import pypdfium2 as pdfium  # type: ignore
    except ImportError:
        return _err("pypdfium2 not installed; cannot read PDF")

    doc = pdfium.PdfDocument(str(path))
    try:
        total_pages = len(doc)
        targets = _parse_pages(pages, total_pages) if pages.strip() else list(range(total_pages))
        total_targets = len(targets)
        chunks: list[str] = []
        ocr_pages = 0
        skipped_pages = 0
        for idx, i in enumerate(targets, start=1):
            if not 0 <= i < total_pages:
                continue
            page = doc[i]
            text = page.get_textpage().get_text_range().strip()
            if len(text) >= _MIN_TEXT_PER_PAGE:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
                emit_progress(
                    "reading_pdf",
                    current=idx,
                    total=total_targets,
                    message=f"page {i + 1}/{total_pages}",
                )
                continue
            if not _ocr_available():
                skipped_pages += 1
                emit_progress(
                    "reading_pdf",
                    current=idx,
                    total=total_targets,
                    message=f"page {i + 1}/{total_pages} (skipped: no OCR)",
                )
                continue
            bitmap = page.render(scale=300 / 72)
            img = bitmap.to_numpy()
            ocr_text = _ocr_image_array(img)
            if ocr_text.strip():
                chunks.append(f"--- Page {i + 1} [OCR] ---\n{ocr_text}")
                ocr_pages += 1
            elif text:
                chunks.append(f"--- Page {i + 1} ---\n{text}")
            emit_progress(
                "reading_pdf",
                current=idx,
                total=total_targets,
                message=f"page {i + 1}/{total_pages} (OCR)" if ocr_text.strip() else f"page {i + 1}/{total_pages}",
            )
        full = "\n\n".join(chunks)
        if not full and skipped_pages > 0:
            engine = _get_ocr()
            hint = get_ocr_install_hint(engine)
            return _err(
                f"All {total_pages} page(s) are scanned/image pages with no "
                f"extractable text, and no OCR engine is available. {hint}"
            )
        return _envelope(
            path, "pdf", full,
            total_pages=total_pages,
            pages_read=len(targets),
            ocr_pages=ocr_pages,
            skipped_pages=skipped_pages,
        )
    finally:
        doc.close()


# ---------------- DOCX ----------------

def _read_docx(path: Path) -> str:
    try:
        import docx  # type: ignore
    except ImportError:
        return _err("python-docx not installed; run: pip install python-docx")

    doc = docx.Document(str(path))
    parts: list[str] = [p.text for p in doc.paragraphs if p.text.strip()]
    for t_idx, table in enumerate(doc.tables, start=1):
        parts.append(f"\n--- Table {t_idx} ---")
        for row in table.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    return _envelope(path, "docx", "\n".join(parts), paragraphs=len(doc.paragraphs), tables=len(doc.tables))


# ---------------- Excel ----------------

def _read_excel(path: Path) -> str:
    try:
        import pandas as pd  # type: ignore
    except ImportError:
        return _err("pandas not installed; cannot read Excel")

    xls = pd.ExcelFile(path)
    parts: list[str] = []
    sheet_info: list[dict[str, Any]] = []
    total_sheets = len(xls.sheet_names)
    for idx, name in enumerate(xls.sheet_names, start=1):
        emit_progress(
            "reading_excel",
            current=idx,
            total=total_sheets,
            message=f"sheet {name}",
        )
        df = xls.parse(name, dtype=str)
        preview = df.head(100).to_string(index=False)
        parts.append(f"--- Sheet: {name} ({len(df)} rows × {len(df.columns)} cols) ---\n{preview}")
        sheet_info.append({"name": name, "rows": len(df), "cols": len(df.columns)})
    return _envelope(path, "excel", "\n\n".join(parts), sheets=sheet_info)


# ---------------- PPTX ----------------

def _read_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return _err("python-pptx not installed; run: pip install python-pptx")

    prs = Presentation(str(path))
    parts: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"--- Slide {idx} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs).strip()
                    if line:
                        parts.append(line)
    return _envelope(path, "pptx", "\n".join(parts), slides=len(prs.slides))


# ---------------- Image OCR ----------------

def _read_image(path: Path) -> str:
    try:
        import numpy as np  # type: ignore
        from PIL import Image  # type: ignore
    except ImportError:
        return _err("Pillow + numpy required for image OCR")

    try:
        img = np.array(Image.open(path).convert("RGB"))
    except Exception as exc:
        return _err(f"Failed to open image: {exc}")

    if not _ocr_available():
        engine = _get_ocr()
        hint = get_ocr_install_hint(engine)
        return _err(
            f"This image requires OCR to extract text, but no OCR engine is "
            f"available. {hint}"
        )

    text = _ocr_image_array(img)
    if not text.strip():
        return _envelope(path, "image", "", note="OCR returned no text (empty or unreadable image)")
    return _envelope(path, "image", text)


# ---------------- Plain text ----------------

def _read_text(path: Path) -> str:
    """Read a text-like file with encoding fallback."""
    data = path.read_bytes()
    last_err: Exception | None = None
    for enc in _ENCODING_FALLBACK:
        try:
            decoded = data.decode(enc)
            return _envelope(path, "text", decoded, encoding=enc, size=len(data))
        except UnicodeDecodeError as exc:
            last_err = exc
    return _err(f"Failed to decode file with any of {_ENCODING_FALLBACK}: {last_err}")


# ---------------- Dispatcher ----------------

_HANDLERS: dict[str, Callable[[Path], str]] = {
    ".docx": _read_docx,
    ".xlsx": _read_excel,
    ".xls": _read_excel,
    ".pptx": _read_pptx,
}


def read_document(
    file_path: str, pages: str = "", *, offset: int = 0, max_chars: int | None = None
) -> str:
    """Read any supported document; dispatch by extension.

    Args:
        file_path: Absolute path to the file.
        pages: Only used for PDF — e.g. "1-10", "5", "1,3,5-8"; empty = all.
        offset: Character offset to start the read window at (for paging).
        max_chars: Read-window size; defaults to ``_DEFAULT_MAX_CHARS`` and is
            clamped to ``_HARD_MAX_CHARS``.

    Returns:
        JSON envelope: status, file, format, char_count, offset, returned,
        next_offset, truncated, text, plus format-specific metadata. When the
        document exceeds the window, follow ``next_offset`` to read the rest.
    """
    try:
        path = safe_document_path(file_path)
    except ValueError as exc:
        return _err(str(exc))
    if not path.exists():
        return _err(f"File not found: {file_path}")
    if not path.is_file():
        return _err(f"Not a file: {file_path}")

    window = (int(offset), _DEFAULT_MAX_CHARS if max_chars is None else int(max_chars))
    token = _window_params.set(window)
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            return _read_pdf(path, pages)
        if ext in _HANDLERS:
            return _HANDLERS[ext](path)
        if ext in _IMAGE_EXTS:
            return _read_image(path)
        if ext in _TEXT_EXTS or ext == "":
            return _read_text(path)
        # Unknown extension: best-effort text read
        return _read_text(path)
    except Exception as exc:
        return _err(f"{type(exc).__name__}: {exc}")
    finally:
        _window_params.reset(token)


class DocReaderTool(BaseTool):
    """Universal document reader — PDF/Word/Excel/PowerPoint/images/text."""

    name = "read_document"
    description = (
        "Read a document of any common format: PDF, Word (.docx), Excel "
        "(.xlsx/.xls), PowerPoint (.pptx), images (OCR), or plain text "
        "(txt/md/json/yaml/csv/html/code files). Returns extracted text in "
        "a unified JSON envelope. For PDFs, accepts an optional `pages` range. "
        "Large documents are paged: follow `next_offset` to read the remainder."
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_path": {"type": "string", "description": "Absolute path to the file."},
            "pages": {
                "type": "string",
                "description": "PDF only: page range (e.g. '1-10', '5', '1,3,5-8'). Ignored for other formats.",
                "default": "",
            },
            "offset": {
                "type": "integer",
                "description": "Character offset to start reading from; pass the previous response's `next_offset` to continue.",
                "default": 0,
            },
            "max_chars": {
                "type": "integer",
                "description": f"Read-window size in characters (default {_DEFAULT_MAX_CHARS}, clamped to {_HARD_MAX_CHARS}).",
            },
        },
        "required": ["file_path"],
    }
    repeatable = True

    def execute(self, **kwargs: Any) -> str:
        raw_max = kwargs.get("max_chars")
        return read_document(
            kwargs["file_path"],
            kwargs.get("pages", ""),
            offset=int(kwargs.get("offset", 0) or 0),
            max_chars=int(raw_max) if raw_max not in (None, "") else None,
        )
